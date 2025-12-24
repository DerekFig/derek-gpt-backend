# main.py

import io
import json
import os
import hashlib
from uuid import UUID
from typing import Optional, List

import fitz  # PyMuPDF
from PIL import Image
from pypdf import PdfReader
from docx import Document as DocxDocument
from io import BytesIO

from pptx import Presentation
from openpyxl import load_workbook

from fastapi import (
    FastAPI,
    Depends,
    UploadFile,
    File,
    Form,
    HTTPException,
    Header,
    Query,
    status,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse

from pydantic import BaseModel

from sqlalchemy import text
from sqlalchemy.orm import Session

from openai import OpenAI, RateLimitError

from db import get_db

import base64

from security import verify_internal_key

from sqlalchemy import text

@app.get("/debug/ingest-tables", dependencies=[Depends(verify_internal_key)])
def debug_ingest_tables(db: Session = Depends(get_db)):
    rows = db.execute(text("""
        SELECT table_name
        FROM information_schema.tables
        WHERE table_schema = 'public'
          AND table_name IN ('ingest_jobs','ingest_job_items')
        ORDER BY table_name;
    """)).fetchall()
    return {"tables": [r.table_name for r in rows]}


# ----------------------------
# NEW: Supabase Storage support
# ----------------------------
import requests
from supabase import create_client, Client


# --------------------------------------------------------
# Environment-backed keys
# --------------------------------------------------------

INTERNAL_BACKEND_KEY = os.getenv("INTERNAL_BACKEND_KEY")

# --------------------------------------------------------
# Admin API key + guard
# --------------------------------------------------------

ADMIN_API_KEY = os.getenv("ADMIN_API_KEY", "changeme-reset-key")


def verify_admin(x_admin_api_key: str = Header(None)):
    """
    Simple admin guard for /admin endpoints.
    Callers must send X-Admin-Api-Key header that matches ADMIN_API_KEY.
    """
    if x_admin_api_key != ADMIN_API_KEY:
        raise HTTPException(status_code=403, detail="Forbidden")


# --------------------------------------------------------
# NEW: Supabase Storage config (backend-only)
# --------------------------------------------------------

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
SUPABASE_STORAGE_BUCKET = os.getenv("SUPABASE_STORAGE_BUCKET", "documents")

supabase: Optional[Client] = None
if SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY:
    supabase = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)


def download_from_supabase_storage(bucket: str, path: str) -> bytes:
    """
    Download object bytes from Supabase Storage via a short-lived signed URL.
    Requires SUPABASE_URL + SUPABASE_SERVICE_ROLE_KEY in backend env.
    """
    if supabase is None:
        raise HTTPException(
            status_code=500,
            detail="Supabase Storage not configured on backend (missing SUPABASE_URL or SUPABASE_SERVICE_ROLE_KEY).",
        )

    signed = supabase.storage.from_(bucket).create_signed_url(path, 60)  # 60 seconds
    signed_url = signed.get("signedURL") or signed.get("signedUrl")
    if not signed_url:
        raise HTTPException(status_code=500, detail="Failed to create signed URL for storage object.")

    r = requests.get(signed_url, timeout=120)
    if r.status_code != 200:
        raise HTTPException(
            status_code=400,
            detail=f"Failed to download file from storage (HTTP {r.status_code}).",
        )

    return r.content


# --------------------------------------------------------
# FastAPI app & CORS
# --------------------------------------------------------

app = FastAPI()

from ingest_routes import router as ingest_router
app.include_router(ingest_router)

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "https://derek-gpt-frontend.vercel.app",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# --------------------------------------------------------
# OCR HELPERS USING OPENAI VISION
# --------------------------------------------------------

vision_client = OpenAI()  # uses OPENAI_API_KEY from env


def _image_to_base64(image: Image.Image) -> str:
    buf = BytesIO()
    image.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def ocr_pages_with_openai(page_images):
    """
    Given a list of PIL Images (one per PDF page),
    call OpenAI Vision to extract ALL visible text from each page.
    Returns a single big string.
    """
    ocr_chunks = []

    for idx, img in enumerate(page_images):
        b64 = _image_to_base64(img)

        resp = vision_client.chat.completions.create(
            model="gpt-4.1",
            temperature=0,
            max_tokens=1200,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "You are performing OCR on a slide or PDF page.\n"
                                "Transcribe **every piece of visible text** exactly as written.\n"
                                "- Include people's names, job titles, company names, and labels.\n"
                                "- Do NOT summarize, rephrase, or skip small text.\n"
                                "- Preserve line breaks where reasonable.\n"
                                "- Output plain text only, no commentary, no bullets added by you."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{b64}"},
                        },
                    ],
                }
            ],
        )

        text_out = (resp.choices[0].message.content or "").strip()
        if text_out:
            ocr_chunks.append(f"[PAGE {idx+1} OCR]\n{text_out}")

    return "\n\n".join(ocr_chunks)


def extract_text_from_pdf_bytes_with_ocr(file_bytes: bytes) -> str:
    """
    Use PyMuPDF to render each page of the PDF (from bytes) to an image,
    then run OCR using OpenAI Vision on each page.
    Returns all OCR text as one string.
    """
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        print(f"PyMuPDF failed to open PDF: {e}")
        return ""

    page_images = []
    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            page_images.append(img)
    except Exception as e:
        print(f"Error rendering PDF pages to images: {e}")
        return ""

    if not page_images:
        return ""

    return ocr_pages_with_openai(page_images)


# --------------------------------------------------------
# Admin reset config
# --------------------------------------------------------

class ResetWorkspaceRequest(BaseModel):
    tenant_id: UUID
    workspace_id: UUID


# --------------------------------------------------------
# OpenAI + embedding model config
# --------------------------------------------------------

EMBEDDING_MODEL = "text-embedding-3-large"
# Must match pgvector: vector(1024)
EMBEDDING_DIM = 1024

client = OpenAI()  # uses OPENAI_API_KEY from env


# --------------------------------------------------------
# Pydantic models
# --------------------------------------------------------

class ChatMessage(BaseModel):
    role: str  # "user" or "assistant"
    content: str


class ChatRequest(BaseModel):
    tenant_id: str
    workspace_id: str
    query: str
    top_k: int = 12
    history: Optional[List[ChatMessage]] = None


class TenantCreate(BaseModel):
    name: str
    type: str = "client"


class WorkspaceCreate(BaseModel):
    tenant_id: str
    name: str


class WorkspaceUpdate(BaseModel):
    system_prompt: Optional[str] = None
    default_model: Optional[str] = None
    temperature: Optional[float] = None


class DocumentIn(BaseModel):
    tenant_id: str
    workspace_id: str
    original_filename: str
    content: str
    metadata: Optional[dict] = None


# ----------------------------
# NEW: ingest-from-storage model
# ----------------------------
class StorageIngestRequest(BaseModel):
    tenant_id: str
    workspace_id: str
    original_filename: str
    storage_path: str
    metadata: Optional[dict] = None


class QueryRequest(BaseModel):
    tenant_id: str
    workspace_id: str
    query: str
    top_k: int = 12


class QueryResult(BaseModel):
    document_id: str
    chunk_index: int
    score: float
    chunk_text: str
    original_filename: Optional[str] = None
    source_type: Optional[str] = None


# --------------------------------------------------------
# Helper functions
# --------------------------------------------------------

def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def get_workspace_config(db: Session, workspace_id: str):
    row = db.execute(
        text("""
            SELECT system_prompt, default_model, temperature
            FROM workspaces
            WHERE id = :wid
        """),
        {"wid": UUID(workspace_id)},
    ).fetchone()

    if not row:
        return {
            "system_prompt": None,
            "default_model": "gpt-4.1-mini",
            "temperature": 0.2,
        }

    return {
        "system_prompt": row.system_prompt,
        "default_model": row.default_model or "gpt-4.1-mini",
        "temperature": float(row.temperature) if row.temperature is not None else 0.2,
    }


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 200) -> List[str]:
    chunks: List[str] = []
    start = 0
    n = len(text)

    while start < n:
        end = min(start + chunk_size, n)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap

    return chunks


SOURCE_TYPES = [
    "deck",
    "email",
    "transcript",
    "notes",
    "document",
    "interview",
    "call",
]


def classify_source_type(filename: str, content: str) -> str:
    lower_name = (filename or "").lower()

    # Heuristics
    if any(ext in lower_name for ext in [".ppt", ".pptx", " deck", " slides"]):
        return "deck"

    if "transcript" in lower_name or "otter" in lower_name:
        return "transcript"
    if any(word in lower_name for word in ["call", "zoom", "meeting"]):
        return "call"

    if "notes" in lower_name:
        return "notes"

    if "interview" in lower_name:
        return "interview"

    if any(ext in lower_name for ext in [".eml", ".msg"]) or "email" in lower_name:
        return "email"

    # Optional LLM classification
    snippet = (content or "")[:2000]

    try:
        completion = client.chat.completions.create(
            model="gpt-4.1-mini",
            temperature=0,
            max_tokens=16,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a classifier. "
                        "Given a filename and text, return exactly ONE of: "
                        "deck, email, transcript, notes, document, interview, call."
                    ),
                },
                {
                    "role": "user",
                    "content": f"Filename: {filename}\n\nText snippet:\n{snippet}",
                },
            ],
        )
        label = (completion.choices[0].message.content or "").strip().lower()
        if label in SOURCE_TYPES:
            return label
        for t in SOURCE_TYPES:
            if t in label:
                return t
    except Exception as e:
        print("WARN: classify_source_type fallback due to error:", repr(e))

    return "document"


def embed_text(text_in: str) -> list[float]:
    try:
        response = client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=text_in,
            dimensions=EMBEDDING_DIM,
        )
        return response.data[0].embedding

    except RateLimitError as e:
        print("OpenAI RateLimitError in embed_text:", repr(e))
        raise HTTPException(
            status_code=503,
            detail="Embedding service unavailable: OpenAI rate limit / quota exceeded.",
        )

    except Exception as e:
        print("Unexpected error in embed_text:", repr(e))
        raise HTTPException(
            status_code=500,
            detail="Unexpected error while creating embedding.",
        )


def embedding_to_pgvector_literal(embedding: list[float]) -> str:
    """
    pgvector expects a text literal like:  '[0.1, -0.2, ...]'  which we then CAST to vector(1024).
    Passing a raw Python list becomes a numeric[] parameter, and Postgres will reject:
      operator does not exist: vector <=> numeric[]
    """
    if not embedding:
        return "[]"
    # Keep full precision; pgvector accepts standard float text.
    return "[" + ",".join(str(float(x)) for x in embedding) + "]"


def ingest_document_text(
    db: Session,
    tenant_id: str,
    workspace_id: str,
    original_filename: str,
    content: str,
    metadata: Optional[dict] = None,
    *,
    document_id: Optional[str] = None,
    file_hash: Optional[str] = None,
):
    metadata = metadata or {}

    source_type = metadata.get("source_type")
    if not source_type:
        source_type = classify_source_type(original_filename, content)

    replaced = False

    # If we are reusing an existing document, clear old embeddings first.
    if document_id:
        replaced = True
        db.execute(
            text("""
                DELETE FROM embeddings
                WHERE document_id = :document_id
                  AND workspace_id = :workspace_id
            """),
            {
                "document_id": UUID(document_id),
                "workspace_id": UUID(workspace_id),
            },
        )

        # Optional: update the document row's content/metadata/hash to match the latest ingest
        db.execute(
            text("""
                UPDATE documents
                SET
                    original_filename = :original_filename,
                    content_text = :content_text,
                    metadata = :metadata,
                    source_type = :source_type,
                    primary_embedding_model = :primary_embedding_model,
                    file_hash = :file_hash
                WHERE id = :document_id
                  AND workspace_id = :workspace_id
            """),
            {
                "document_id": UUID(document_id),
                "workspace_id": UUID(workspace_id),
                "original_filename": original_filename,
                "content_text": content,
                "metadata": json.dumps(metadata) if metadata else None,
                "source_type": source_type,
                "primary_embedding_model": EMBEDDING_MODEL,
                "file_hash": file_hash,
            },
        )
    else:
        # New document insert (now includes file_hash)
        result = db.execute(
            text("""
                INSERT INTO documents (
                    workspace_id,
                    original_filename,
                    content_text,
                    metadata,
                    source_type,
                    primary_embedding_model,
                    file_hash
                )
                VALUES (
                    :workspace_id,
                    :original_filename,
                    :content_text,
                    :metadata,
                    :source_type,
                    :primary_embedding_model,
                    :file_hash
                )
                RETURNING id;
            """),
            {
                "workspace_id": UUID(workspace_id),
                "original_filename": original_filename,
                "content_text": content,
                "metadata": json.dumps(metadata) if metadata else None,
                "source_type": source_type,
                "primary_embedding_model": EMBEDDING_MODEL,
                "file_hash": file_hash,
            },
        )
        doc_row = result.fetchone()
        document_id = str(doc_row.id)

    chunks = chunk_text(content)

    for idx, chunk in enumerate(chunks):
        embedding = embed_text(chunk)

        if isinstance(embedding, str):
            raise RuntimeError(
                f"embed_text returned a string instead of a vector for chunk {idx}"
            )

        db.execute(
            text("""
                INSERT INTO embeddings (
                    tenant_id,
                    workspace_id,
                    document_id,
                    chunk_index,
                    chunk_text,
                    embedding,
                    embedding_model
                )
                VALUES (
                    :tenant_id,
                    :workspace_id,
                    :document_id,
                    :chunk_index,
                    :chunk_text,
                    :embedding,
                    :embedding_model
                );
            """),
            {
                "tenant_id": UUID(tenant_id),
                "workspace_id": UUID(workspace_id),
                "document_id": UUID(document_id),
                "chunk_index": idx,
                "chunk_text": chunk,
                "embedding": embedding,              # ✅ vector values
                "embedding_model": EMBEDDING_MODEL,  # ✅ model name
            },
        )

    db.commit()

    return {
        "document_id": str(document_id),
        "num_chunks": len(chunks),
        "source_type": source_type,
        "replaced": replaced,
        "file_hash": file_hash,
    }


def extract_text_from_pptx_bytes(content_bytes: bytes) -> str:
    if not content_bytes:
        return ""

    prs = Presentation(BytesIO(content_bytes))
    texts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_text_parts.append(t)

        if slide_text_parts:
            texts.append(f"[Slide {slide_idx}]")
            texts.append("\n".join(slide_text_parts))

    return "\n\n".join(texts).strip()


def extract_text_from_xlsx_bytes(content_bytes: bytes) -> str:
    if not content_bytes:
        return ""

    wb = load_workbook(BytesIO(content_bytes), data_only=True)
    lines: list[str] = []

    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) for cell in row if cell is not None]
            if row_values:
                lines.append("\t".join(row_values))
        lines.append("")

    return "\n".join(lines).strip()


def extract_text_from_upload(file: UploadFile) -> str:
    filename = file.filename or "unnamed"
    lower_name = filename.lower()

    content_bytes = file.file.read()

    if lower_name.endswith(".txt") or lower_name.endswith(".md"):
        return content_bytes.decode("utf-8", errors="ignore")

    if lower_name.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(content_bytes))
        pages_text = []
        for page in reader.pages:
            pages_text.append(page.extract_text() or "")
        return "\n".join(pages_text)

    if lower_name.endswith(".docx"):
        doc = DocxDocument(io.BytesIO(content_bytes))
        paragraphs = [p.text for p in doc.paragraphs]
        return "\n".join(paragraphs)

    if lower_name.endswith(".pptx"):
        return extract_text_from_pptx_bytes(content_bytes)

    if lower_name.endswith(".ppt"):
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported legacy PowerPoint file '{filename}'. "
                "Please resave the file as .pptx and upload again."
            ),
        )

    if lower_name.endswith(".xlsx"):
        return extract_text_from_xlsx_bytes(content_bytes)

    if lower_name.endswith(".xls"):
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported legacy Excel file '{filename}'. "
                "Please resave the file as .xlsx and upload again."
            ),
        )

    raise HTTPException(
        status_code=400,
        detail=(
            f"Unsupported file type for '{filename}'. "
            "Supported: .txt, .md, .pdf, .docx, .pptx, .xlsx"
        ),
    )


# --------------------------------------------------------
# Routes
# --------------------------------------------------------

@app.get("/health")
def health():
    return {"status": "ok"}


@app.post(
    "/admin/reset-workspace",
    dependencies=[Depends(verify_internal_key)],
)
def reset_workspace(
    payload: ResetWorkspaceRequest,
    db: Session = Depends(get_db),
    _: None = Depends(verify_admin),
):
    tenant_id = str(payload.tenant_id)
    workspace_id = str(payload.workspace_id)

    workspace_row = db.execute(
        text("""
            SELECT id, tenant_id
            FROM workspaces
            WHERE id = :workspace_id
            LIMIT 1
        """),
        {"workspace_id": workspace_id},
    ).fetchone()

    if workspace_row is None:
        raise HTTPException(status_code=404, detail="Workspace not found")

    if str(workspace_row.tenant_id) != tenant_id:
        print(
            "WARNING: reset_workspace called with tenant_id",
            tenant_id,
            "but workspace actually belongs to",
            str(workspace_row.tenant_id),
        )

    try:
        db.execute(
            text("""
                DELETE FROM embeddings
                WHERE workspace_id = :workspace_id
            """),
            {"workspace_id": workspace_id},
        )

        db.execute(
            text("""
                DELETE FROM documents
                WHERE workspace_id = :workspace_id
            """),
            {"workspace_id": workspace_id},
        )

        db.commit()
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to reset workspace: {e}")

    return {
        "status": "ok",
        "message": "Workspace reset. Documents and embeddings deleted.",
        "tenant_id": tenant_id,
        "workspace_id": workspace_id,
    }


@app.get("/tenants", dependencies=[Depends(verify_internal_key)])
def list_tenants(db: Session = Depends(get_db)):
    result = db.execute(
        text("SELECT id, name, type, created_at FROM tenants ORDER BY created_at ASC;")
    )
    rows = result.fetchall()

    return [
        {
            "id": str(row.id),
            "name": row.name,
            "type": row.type,
            "created_at": row.created_at,
        }
        for row in rows
    ]


@app.post("/tenants", dependencies=[Depends(verify_internal_key)])
def create_tenant(tenant: TenantCreate, db: Session = Depends(get_db)):
    stmt = text("""
        INSERT INTO tenants (name, type)
        VALUES (:name, :type)
        RETURNING id, name, type, created_at
    """)

    result = db.execute(
        stmt,
        {"name": tenant.name, "type": tenant.type},
    )

    row = result.fetchone()
    db.commit()

    return {
        "id": str(row.id),
        "name": row.name,
        "type": row.type,
        "created_at": row.created_at.isoformat(),
    }


@app.get("/workspaces", dependencies=[Depends(verify_internal_key)])
def list_workspaces(tenant_id: str, db: Session = Depends(get_db)):
    stmt = text("""
        SELECT id, tenant_id, name, created_at
        FROM workspaces
        WHERE tenant_id = :tenant_id
        ORDER BY created_at ASC
    """)

    rows = db.execute(stmt, {"tenant_id": tenant_id}).fetchall()

    return [
        {
            "id": str(row.id),
            "tenant_id": str(row.tenant_id),
            "name": row.name,
            "created_at": row.created_at.isoformat(),
        }
        for row in rows
    ]


@app.post("/workspaces", dependencies=[Depends(verify_internal_key)])
def create_workspace(payload: WorkspaceCreate, db: Session = Depends(get_db)):
    stmt = text("""
        INSERT INTO workspaces (tenant_id, name)
        VALUES (:tenant_id, :name)
        RETURNING id, tenant_id, name, created_at
    """)

    result = db.execute(
        stmt,
        {
            "tenant_id": str(payload.tenant_id),
            "name": payload.name,
        },
    )

    row = result.fetchone()
    db.commit()

    return {
        "id": str(row.id),
        "tenant_id": str(row.tenant_id),
        "name": row.name,
        "created_at": row.created_at.isoformat(),
    }


@app.get("/documents", dependencies=[Depends(verify_internal_key)])
def list_documents(
    tenant_id: str,
    workspace_id: str,
    db: Session = Depends(get_db),
):
    sql = text("""
        SELECT
            d.id,
            d.original_filename,
            d.created_at,
            d.source_type,
            COALESCE(c.chunk_count, 0) AS chunk_count
        FROM documents d
        LEFT JOIN (
            SELECT document_id, COUNT(*) AS chunk_count
            FROM embeddings
            GROUP BY document_id
        ) c ON c.document_id = d.id
        WHERE d.workspace_id = :workspace_id
        ORDER BY d.created_at DESC;
    """)

    rows = db.execute(
        sql,
        {"workspace_id": UUID(workspace_id)},
    ).fetchall()

    return [
        {
            "id": str(row.id),
            "original_filename": row.original_filename,
            "created_at": row.created_at.isoformat() if row.created_at else None,
            "chunk_count": int(row.chunk_count),
            "source_type": row.source_type,
        }
        for row in rows
    ]


@app.delete(
    "/documents/{document_id}",
    status_code=status.HTTP_204_NO_CONTENT,
    dependencies=[Depends(verify_internal_key)],
)
def delete_document(
    document_id: str,
    tenant_id: str = Query(...),
    workspace_id: str = Query(...),
    db: Session = Depends(get_db),
):
    check_sql = text("""
        SELECT id
        FROM documents
        WHERE id = :document_id
          AND workspace_id = :workspace_id
        LIMIT 1
    """)

    doc_row = db.execute(
        check_sql,
        {
            "document_id": UUID(document_id),
            "workspace_id": UUID(workspace_id),
        },
    ).fetchone()

    if doc_row is None:
        raise HTTPException(status_code=404, detail="Document not found in workspace")

    try:
        db.execute(
            text("""
                DELETE FROM embeddings
                WHERE document_id = :document_id
                  AND workspace_id = :workspace_id
                  AND tenant_id = :tenant_id
            """),
            {
                "document_id": UUID(document_id),
                "workspace_id": UUID(workspace_id),
                "tenant_id": UUID(tenant_id),
            },
        )

        db.execute(
            text("""
                DELETE FROM documents
                WHERE id = :document_id
                  AND workspace_id = :workspace_id
            """),
            {
                "document_id": UUID(document_id),
                "workspace_id": UUID(workspace_id),
            },
        )

        db.commit()
    except Exception as e:
        db.rollback()
        raise HTTPException(
            status_code=500,
            detail=f"Failed to delete document: {e}",
        )

    return


@app.get("/debug/db", dependencies=[Depends(verify_internal_key)])
def debug_db(db: Session = Depends(get_db)):
    result = db.execute(text("SELECT current_database(), current_user, current_schema();"))
    db_name, user, schema = result.fetchone()

    ws_result = db.execute(text("SELECT id, tenant_id, name FROM workspaces"))
    workspaces = [
        {"id": str(r.id), "tenant_id": str(r.tenant_id), "name": r.name}
        for r in ws_result.fetchall()
    ]

    return {
        "database": db_name,
        "user": user,
        "schema": schema,
        "workspaces": workspaces,
    }


@app.get("/debug/openai", dependencies=[Depends(verify_internal_key)])
def debug_openai():
    try:
        response = client.embeddings.create(
            model="text-embedding-3-small",
            input="ping",
        )
        dim = len(response.data[0].embedding)
        return {
            "status": "ok",
            "model": "text-embedding-3-small",
            "dimension": dim,
        }

    except RateLimitError:
        raise HTTPException(
            status_code=503,
            detail="OpenAI quota / billing issue: please add credits or update your plan.",
        )

    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Unexpected OpenAI error: {e!r}",
        )


@app.post("/query", dependencies=[Depends(verify_internal_key)])
def query_embeddings(
    payload: QueryRequest,
    db: Session = Depends(get_db),
):
    query_embedding = embed_text(payload.query)
    query_vec = embedding_to_pgvector_literal(query_embedding)

    sql = text(f"""
        SELECT
            e.document_id,
            e.chunk_index,
            e.chunk_text,
            1 - (e.embedding <=> CAST(:query_embedding AS vector({EMBEDDING_DIM}))) AS score,
            d.original_filename,
            d.source_type
        FROM embeddings e
        JOIN documents d ON d.id = e.document_id
        WHERE e.tenant_id = :tenant_id
          AND e.workspace_id = :workspace_id
          AND e.embedding_model = :embedding_model
        ORDER BY score DESC
        LIMIT :top_k;
    """)

    rows = db.execute(
        sql,
        {
            "tenant_id": payload.tenant_id,
            "workspace_id": payload.workspace_id,
            "embedding_model": EMBEDDING_MODEL,
            "query_embedding": query_vec,
            "top_k": payload.top_k,
        },
    ).fetchall()

    results: List[QueryResult] = []
    for row in rows:
        results.append(
            QueryResult(
                document_id=str(row.document_id),
                chunk_index=row.chunk_index,
                score=float(row.score),
                chunk_text=row.chunk_text,
                original_filename=row.original_filename,
                source_type=row.source_type,
            )
        )

    return {
        "query": payload.query,
        "results": [r.dict() for r in results],
    }


@app.post("/documents", dependencies=[Depends(verify_internal_key)])
def ingest_document(
    payload: DocumentIn,
    db: Session = Depends(get_db),
):
    ws_sql = text("""
        SELECT id, tenant_id, name
        FROM workspaces
        WHERE id = :wid
    """)

    ws_row = db.execute(ws_sql, {"wid": str(payload.workspace_id)}).fetchone()

    if ws_row is None:
        raise HTTPException(
            status_code=400,
            detail=f"Workspace {payload.workspace_id} not found (pre-insert check)",
        )

    return ingest_document_text(
        db=db,
        tenant_id=payload.tenant_id,
        workspace_id=payload.workspace_id,
        original_filename=payload.original_filename,
        content=payload.content,
        metadata=payload.metadata,
        file_hash=None,
    )


@app.post("/upload-file", dependencies=[Depends(verify_internal_key)])
def upload_file(
    tenant_id: str = Form(...),
    workspace_id: str = Form(...),
    metadata: Optional[str] = Form(None),
    file: List[UploadFile] = File(...),
    db: Session = Depends(get_db),
):
    metadata_dict = None
    if metadata:
        try:
            metadata_dict = json.loads(metadata)
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="metadata must be valid JSON")

    if not file or len(file) == 0:
        raise HTTPException(status_code=400, detail="No files uploaded.")

    results: list = []

    for f in file:
        filename_lower = (f.filename or "").lower()
        text_content: Optional[str] = None

        if filename_lower.endswith(".pdf"):
            file_bytes = f.file.read()
            f.file.seek(0)

            ocr_text = extract_text_from_pdf_bytes_with_ocr(file_bytes)
            print("DEBUG: OCR text length for", f.filename, "=", len(ocr_text or ""))

            if ocr_text and ocr_text.strip():
                text_content = "[USING_OCR]\n" + ocr_text
            else:
                text_content = extract_text_from_upload(f)
        else:
            text_content = extract_text_from_upload(f)

        if not text_content or not text_content.strip():
            continue

        ingest_result = ingest_document_text(
            db=db,
            tenant_id=tenant_id,
            workspace_id=workspace_id,
            original_filename=f.filename or "uploaded_file",
            content=text_content,
            metadata=metadata_dict,
            file_hash=None,
        )
        results.append(ingest_result)

    if not results:
        raise HTTPException(
            status_code=400,
            detail="No text could be extracted from any file.",
        )

    return {"results": results}


# ----------------------------
# NEW: Ingest from Supabase Storage (by storage_path)
# ----------------------------
@app.post("/ingest-from-storage", dependencies=[Depends(verify_internal_key)])
def ingest_from_storage(
    payload: StorageIngestRequest,
    db: Session = Depends(get_db),
):
    file_bytes = download_from_supabase_storage(SUPABASE_STORAGE_BUCKET, payload.storage_path)
    file_hash = sha256_bytes(file_bytes)

    # Check for duplicate in same workspace by (workspace_id, file_hash)
    existing = db.execute(
        text("""
            SELECT id
            FROM documents
            WHERE workspace_id = :workspace_id
              AND file_hash = :file_hash
            LIMIT 1
        """),
        {
            "workspace_id": UUID(payload.workspace_id),
            "file_hash": file_hash,
        },
    ).fetchone()
    existing_document_id = str(existing.id) if existing else None

    filename_lower = (payload.original_filename or "").lower()
    text_content: Optional[str] = None

    if filename_lower.endswith(".pdf"):
        ocr_text = extract_text_from_pdf_bytes_with_ocr(file_bytes)
        print("DEBUG: OCR text length for", payload.original_filename, "=", len(ocr_text or ""))

        if ocr_text and ocr_text.strip():
            text_content = "[USING_OCR]\n" + ocr_text
        else:
            reader = PdfReader(io.BytesIO(file_bytes))
            pages_text = [(p.extract_text() or "") for p in reader.pages]
            text_content = "\n".join(pages_text)

    elif filename_lower.endswith(".docx"):
        doc = DocxDocument(io.BytesIO(file_bytes))
        text_content = "\n".join([p.text for p in doc.paragraphs])

    elif filename_lower.endswith(".pptx"):
        text_content = extract_text_from_pptx_bytes(file_bytes)

    elif filename_lower.endswith(".xlsx"):
        text_content = extract_text_from_xlsx_bytes(file_bytes)

    elif filename_lower.endswith(".txt") or filename_lower.endswith(".md"):
        text_content = file_bytes.decode("utf-8", errors="ignore")

    elif filename_lower.endswith(".ppt"):
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported legacy PowerPoint file '{payload.original_filename}'. "
                "Please resave the file as .pptx and upload again."
            ),
        )

    elif filename_lower.endswith(".xls"):
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported legacy Excel file '{payload.original_filename}'. "
                "Please resave the file as .xlsx and upload again."
            ),
        )

    else:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported file type for '{payload.original_filename}'. "
                "Supported: .txt, .md, .pdf, .docx, .pptx, .xlsx"
            ),
        )

    if not text_content or not text_content.strip():
        raise HTTPException(status_code=400, detail="No text could be extracted from the file.")

    ingest_result = ingest_document_text(
        db=db,
        tenant_id=payload.tenant_id,
        workspace_id=payload.workspace_id,
        original_filename=payload.original_filename,
        content=text_content,
        metadata=payload.metadata,
        document_id=existing_document_id,
        file_hash=file_hash,
    )

    return {"result": ingest_result}


@app.post(
    "/chat",
    dependencies=[Depends(verify_internal_key)],
)
def chat_with_workspace(
    payload: ChatRequest,
    db: Session = Depends(get_db),
):
    """
    RAG-style multi-turn chat (non-streaming).
    """
    query_embedding = embed_text(payload.query)
    query_vec = embedding_to_pgvector_literal(query_embedding)

    # UPDATED: join documents so we can include filename/source_type in context
    sql = text(f"""
        SELECT
            e.document_id,
            e.chunk_index,
            e.chunk_text,
            1 - (e.embedding <=> CAST(:query_embedding AS vector({EMBEDDING_DIM}))) AS score,
            d.original_filename,
            d.source_type
        FROM embeddings e
        JOIN documents d ON d.id = e.document_id
        WHERE e.tenant_id = :tenant_id
          AND e.workspace_id = :workspace_id
          AND e.embedding_model = :embedding_model
        ORDER BY score DESC
        LIMIT :top_k;
    """)

    rows = db.execute(
        sql,
        {
            "tenant_id": payload.tenant_id,
            "workspace_id": payload.workspace_id,
            "embedding_model": EMBEDDING_MODEL,
            "query_embedding": query_vec,
            "top_k": payload.top_k,
        },
    ).fetchall()

    # UPDATED: richer context with provenance
    contexts: List[str] = []
    for row in rows:
        fname = getattr(row, "original_filename", None) or "unknown"
        stype = getattr(row, "source_type", None) or "unknown"
        contexts.append(
            f"[SOURCE: {fname} | type={stype} | doc={row.document_id} | chunk={row.chunk_index} | score={float(row.score):.3f}]\n"
            f"{row.chunk_text}"
        )

    if not contexts:
        return {
            "answer": "I couldn't find any relevant documents in this workspace yet.",
            "results": [],
        }

    history_lines: List[str] = []
    if payload.history:
        recent = payload.history[-10:]
        for m in recent:
            role = m.role.lower()
            if role not in ("user", "assistant"):
                continue
            tag = "User" if role == "user" else "Assistant"
            history_lines.append(f"{tag}: {m.content.strip()}")

    history_block = "\n".join(history_lines) if history_lines else "None so far."

    q = (payload.query or "").strip()
    q_lower = q.lower()

    interpretation_hint = ""
    if q_lower.startswith("who is ") or q_lower.startswith("who are "):
        interpretation_hint = (
            "Interpret 'who is/are X' flexibly. If X appears to be a company or organization, "
            "answer what it is and identify the principals/founders/leaders mentioned in the context.\n\n"
        )

    # UPDATED: extraction-first prompt reduces overly literal failures
    base_context_prompt = (
        interpretation_hint
        + "You are answering using only the provided document context.\n\n"
          "DOCUMENT CONTEXT:\n"
        + "\n\n---\n\n".join(contexts)
        + "\n\nCONVERSATION HISTORY:\n"
        + history_block
        + "\n\nUSER QUESTION:\n"
        + q
        + "\n\nINSTRUCTIONS:\n"
          "1) First, extract the 3–8 most relevant facts from the context as short bullets.\n"
          "2) Then answer the question in 2–6 sentences using those facts.\n"
          "3) If the context truly contains no relevant facts, say: 'I don't have that in the documents yet.'\n"
    )

    cfg = get_workspace_config(db, payload.workspace_id)
    system_prompt = (
        cfg["system_prompt"]
        or "You are DerekGPT, a helpful assistant answering questions based on the provided context."
    )
    model_name = cfg["default_model"]
    temperature = cfg["temperature"]

    completion = client.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": base_context_prompt},
        ],
        temperature=temperature,
    )

    answer = completion.choices[0].message.content

    # UPDATED: return provenance to the frontend for debugging
    results = []
    for row in rows:
        results.append(
            {
                "document_id": str(row.document_id),
                "chunk_index": row.chunk_index,
                "score": float(row.score),
                "chunk_text": row.chunk_text,
                "original_filename": getattr(row, "original_filename", None),
                "source_type": getattr(row, "source_type", None),
            }
        )

    return {"answer": answer, "results": results}


@app.post("/chat-stream", dependencies=[Depends(verify_internal_key)])
def chat_with_workspace_stream(
    payload: ChatRequest,
    db: Session = Depends(get_db),
):
    """
    Streaming RAG chat endpoint.
    """
    query_embedding = embed_text(payload.query)
    query_vec = embedding_to_pgvector_literal(query_embedding)

    # UPDATED: join documents so we can include filename/source_type in context
    sql = text(f"""
        SELECT
            e.document_id,
            e.chunk_index,
            e.chunk_text,
            1 - (e.embedding <=> CAST(:query_embedding AS vector({EMBEDDING_DIM}))) AS score,
            d.original_filename,
            d.source_type
        FROM embeddings e
        JOIN documents d ON d.id = e.document_id
        WHERE e.tenant_id = :tenant_id
          AND e.workspace_id = :workspace_id
          AND e.embedding_model = :embedding_model
        ORDER BY score DESC
        LIMIT :top_k;
    """)

    rows = db.execute(
        sql,
        {
            "tenant_id": payload.tenant_id,
            "workspace_id": payload.workspace_id,
            "embedding_model": EMBEDDING_MODEL,
            "query_embedding": query_vec,
            "top_k": payload.top_k,
        },
    ).fetchall()

    # UPDATED: richer context with provenance
    contexts: List[str] = []
    for row in rows:
        fname = getattr(row, "original_filename", None) or "unknown"
        stype = getattr(row, "source_type", None) or "unknown"
        contexts.append(
            f"[SOURCE: {fname} | type={stype} | doc={row.document_id} | chunk={row.chunk_index} | score={float(row.score):.3f}]\n"
            f"{row.chunk_text}"
        )

    if not contexts:
        def empty_gen():
            yield "I couldn't find any relevant documents in this workspace yet."
        return StreamingResponse(empty_gen(), media_type="text/plain")

    history_lines: List[str] = []
    if payload.history:
        for m in payload.history[-10:]:
            role = m.role.lower()
            if role not in ("user", "assistant"):
                continue
            tag = "User" if role == "user" else "Assistant"
            history_lines.append(f"{tag}: {m.content.strip()}")

    history_block = "\n".join(history_lines) if history_lines else "None so far."

    q = (payload.query or "").strip()
    q_lower = q.lower()

    interpretation_hint = ""
    if q_lower.startswith("who is ") or q_lower.startswith("who are "):
        interpretation_hint = (
            "Interpret 'who is/are X' flexibly. If X appears to be a company or organization, "
            "answer what it is and identify the principals/founders/leaders mentioned in the context.\n\n"
        )

    base_context_prompt = (
        interpretation_hint
        + "You are answering using only the provided document context.\n\n"
          "DOCUMENT CONTEXT:\n"
        + "\n\n---\n\n".join(contexts)
        + "\n\nCONVERSATION HISTORY:\n"
        + history_block
        + "\n\nUSER QUESTION:\n"
        + q
        + "\n\nINSTRUCTIONS:\n"
          "1) First, extract the 3–8 most relevant facts from the context as short bullets.\n"
          "2) Then answer the question in 2–6 sentences using those facts.\n"
          "3) If the context truly contains no relevant facts, say: 'I don't have that in the documents yet.'\n"
    )

    cfg = get_workspace_config(db, payload.workspace_id)
    system_prompt = (
        cfg["system_prompt"]
        or "You are DerekGPT, a helpful assistant answering questions based on the provided context."
    )
    model_name = cfg["default_model"]
    temperature = cfg["temperature"]

    def token_stream():
        try:
            stream = client.chat.completions.create(
                model=model_name,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": base_context_prompt},
                ],
                temperature=temperature,
                stream=True,
            )
            for chunk in stream:
                delta = chunk.choices[0].delta
                content = getattr(delta, "content", None) or ""
                if content:
                    yield content
        except Exception as e:
            print("Streaming error:", repr(e))
            yield "\n[Error while streaming answer]"

    return StreamingResponse(token_stream(), media_type="text/plain")


@app.patch(
    "/workspaces/{workspace_id}",
    dependencies=[Depends(verify_internal_key)],
)
def update_workspace(
    workspace_id: str,
    payload: WorkspaceUpdate,
    db: Session = Depends(get_db),
):
    fields = []
    params = {"workspace_id": UUID(workspace_id)}

    if payload.system_prompt is not None:
        fields.append("system_prompt = :system_prompt")
        params["system_prompt"] = payload.system_prompt

    if payload.default_model is not None:
        fields.append("default_model = :default_model")
        params["default_model"] = payload.default_model

    if payload.temperature is not None:
        fields.append("temperature = :temperature")
        params["temperature"] = payload.temperature

    if not fields:
        return {"status": "ok"}

    sql = text(f"""
        UPDATE workspaces
        SET {", ".join(fields)}
        WHERE id = :workspace_id
        RETURNING id;
    """)

    row = db.execute(sql, params).fetchone()
    db.commit()

    if not row:
        raise HTTPException(status_code=404, detail="Workspace not found")

    return {"status": "ok", "workspace_id": str(row.id)}


@app.get("/whoami")
def whoami():
    return {"running_file": "MAIN_PY_IS_RUNNING_V3"}
