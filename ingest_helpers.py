# ingest_helpers.py

from __future__ import annotations

import base64
import hashlib
import io
import json
import os
from io import BytesIO
from typing import Optional, List
from uuid import UUID
import logging

import fitz  # PyMuPDF
import requests
from PIL import Image
from docx import Document as DocxDocument
from openpyxl import load_workbook
from openai import OpenAI, RateLimitError
from pptx import Presentation
from pypdf import PdfReader

from fastapi import HTTPException, UploadFile
from sqlalchemy import text
from sqlalchemy.orm import Session
from supabase import create_client, Client


# --------------------------------------------------------
# Environment + Supabase client
# --------------------------------------------------------

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
SUPABASE_STORAGE_BUCKET = os.getenv("SUPABASE_STORAGE_BUCKET", "documents")

supabase: Optional[Client] = None
if SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY:
    supabase = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)


def download_from_supabase_storage(bucket: str, path: str) -> bytes:
    """
    Download object bytes from Supabase Storage via a short-lived signed URL.
    Requires SUPABASE_URL + SUPABASE_SERVICE_ROLE_KEY in backend env.
    """
    if supabase is None:
        raise HTTPException(
            status_code=500,
            detail="Supabase Storage not configured on backend (missing SUPABASE_URL or SUPABASE_SERVICE_ROLE_KEY).",
        )

    signed = supabase.storage.from_(bucket).create_signed_url(path, 60)  # 60 seconds
    signed_url = signed.get("signedURL") or signed.get("signedUrl")
    if not signed_url:
        raise HTTPException(status_code=500, detail="Failed to create signed URL for storage object.")

    r = requests.get(signed_url, timeout=120)
    if r.status_code != 200:
        raise HTTPException(
            status_code=400,
            detail=f"Failed to download file from storage (HTTP {r.status_code}).",
        )

    return r.content


# --------------------------------------------------------
# OCR helpers using OpenAI Vision
# --------------------------------------------------------

vision_client = OpenAI()  # uses OPENAI_API_KEY from env


def _image_to_base64(image: Image.Image) -> str:
    buf = BytesIO()
    image.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def ocr_pages_with_openai(page_images):
    """
    Given a list of PIL Images (one per PDF page),
    call OpenAI Vision to extract ALL visible text from each page.
    Returns a single big string.
    """
    ocr_chunks = []

    for idx, img in enumerate(page_images):
        b64 = _image_to_base64(img)

        resp = vision_client.chat.completions.create(
            model="gpt-4.1",
            temperature=0,
            max_tokens=1200,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "You are performing OCR on a slide or PDF page.\n"
                                "Transcribe **every piece of visible text** exactly as written.\n"
                                "- Include people's names, job titles, company names, and labels.\n"
                                "- Do NOT summarize, rephrase, or skip small text.\n"
                                "- Preserve line breaks where reasonable.\n"
                                "- Output plain text only, no commentary, no bullets added by you."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{b64}"},
                        },
                    ],
                }
            ],
        )

        text_out = (resp.choices[0].message.content or "").strip()
        if text_out:
            ocr_chunks.append(f"[PAGE {idx+1} OCR]\n{text_out}")

    return "\n\n".join(ocr_chunks)


def extract_text_from_pdf_bytes_with_ocr(file_bytes: bytes) -> str:
    """
    Use PyMuPDF to render each page of the PDF (from bytes) to an image,
    then run OCR using OpenAI Vision on each page.
    Returns all OCR text as one string.
    """
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        print(f"PyMuPDF failed to open PDF: {e}")
        return ""

    page_images = []
    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            page_images.append(img)
    except Exception as e:
        print(f"Error rendering PDF pages to images: {e}")
        return ""

    if not page_images:
        return ""

    return ocr_pages_with_openai(page_images)


# --------------------------------------------------------
# OpenAI + embedding model config
# --------------------------------------------------------

EMBEDDING_MODEL = "text-embedding-3-large"
EMBEDDING_DIM = 1024  # Must match pgvector: vector(1024)

client = OpenAI()  # uses OPENAI_API_KEY from env


# --------------------------------------------------------
# General helper functions used by ingestion
# --------------------------------------------------------

def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def chunk_text(text_in: str, chunk_size: int = 800, overlap: int = 200) -> List[str]:
    chunks: List[str] = []
    start = 0
    n = len(text_in)

    while start < n:
        end = min(start + chunk_size, n)
        chunk = text_in[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap

    return chunks


SOURCE_TYPES = ["deck", "email", "transcript", "notes", "document", "interview", "call"]


def classify_source_type(filename: str, content: str) -> str:
    lower_name = (filename or "").lower()

    if any(ext in lower_name for ext in [".ppt", ".pptx", " deck", " slides"]):
        return "deck"

    if "transcript" in lower_name or "otter" in lower_name:
        return "transcript"
    if any(word in lower_name for word in ["call", "zoom", "meeting"]):
        return "call"
    if "notes" in lower_name:
        return "notes"
    if "interview" in lower_name:
        return "interview"
    if any(ext in lower_name for ext in [".eml", ".msg"]) or "email" in lower_name:
        return "email"

    snippet = (content or "")[:2000]
    try:
        completion = client.chat.completions.create(
            model="gpt-4.1-mini",
            temperature=0,
            max_tokens=16,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a classifier. "
                        "Given a filename and text, return exactly ONE of: "
                        "deck, email, transcript, notes, document, interview, call."
                    ),
                },
                {"role": "user", "content": f"Filename: {filename}\n\nText snippet:\n{snippet}"},
            ],
        )
        label = (completion.choices[0].message.content or "").strip().lower()
        if label in SOURCE_TYPES:
            return label
        for t in SOURCE_TYPES:
            if t in label:
                return t
    except Exception as e:
        print("WARN: classify_source_type fallback due to error:", repr(e))

    return "document"


def embed_text(text_in: str) -> list[float]:
    try:
        response = client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=text_in,
            dimensions=EMBEDDING_DIM,
        )
        return response.data[0].embedding

    except RateLimitError as e:
        print("OpenAI RateLimitError in embed_text:", repr(e))
        raise HTTPException(
            status_code=503,
            detail="Embedding service unavailable: OpenAI rate limit / quota exceeded.",
        )

    except Exception as e:
        print("Unexpected error in embed_text:", repr(e))
        raise HTTPException(
            status_code=500,
            detail="Unexpected error while creating embedding.",
        )


def embedding_to_pgvector_literal(embedding: list[float]) -> str:
    if not embedding:
        return "[]"
    return "[" + ",".join(str(float(x)) for x in embedding) + "]"


logger = logging.getLogger(__name__)


def ingest_document_text(
    db: Session,
    tenant_id: str,
    workspace_id: str,
    original_filename: str,
    content: str,
    metadata: Optional[dict] = None,
    *,
    document_id: Optional[str] = None,
    file_hash: Optional[str] = None,
):
    """
    Ingests a document into `documents` and its chunks into `embeddings`.

    DEDUPE BEHAVIOR (strict):
      - Requires a UNIQUE index in Postgres:
          create unique index if not exists documents_dedupe_idx
          on documents (workspace_id, file_hash)
          where file_hash is not null;
      - If `document_id` is not provided and (workspace_id, file_hash) already exists,
        this function will SKIP creating a new documents row and SKIP embeddings.
    """
    metadata = metadata or {}

    # truth/debug (safe)
    if file_hash is None:
        logger.error("ingest_document_text called with file_hash=None for %s", original_filename)
    else:
        logger.info("ingest_document_text file_hash=%s filename=%s", file_hash, original_filename)

    # ----------------------------
    # FIX: normalize file_hash so empty/whitespace never routes to NULL insert accidentally
    # ----------------------------
    if isinstance(file_hash, str):
        file_hash = file_hash.strip()
        if file_hash == "":
            file_hash = None

    # Optional sanity logging (does not change behavior)
    if file_hash is not None and isinstance(file_hash, str) and len(file_hash) != 64:
        logger.error("Unexpected file_hash length (expected 64). filename=%s file_hash=%r", original_filename, file_hash)

    source_type = metadata.get("source_type")
    if not source_type:
        source_type = classify_source_type(original_filename, content)

    replaced = False

    # ----------------------------
    # Replace existing (explicit id)
    # ----------------------------
    if document_id:
        replaced = True

        db.execute(
            text("""
                DELETE FROM embeddings
                WHERE document_id = :document_id
                  AND workspace_id = :workspace_id
            """),
            {"document_id": UUID(document_id), "workspace_id": UUID(workspace_id)},
        )

        db.execute(
            text("""
                UPDATE documents
                SET
                    original_filename = :original_filename,
                    content_text = :content_text,
                    metadata = :metadata,
                    source_type = :source_type,
                    primary_embedding_model = :primary_embedding_model,
                    file_hash = :file_hash
                WHERE id = :document_id
                  AND workspace_id = :workspace_id
            """),
            {
                "document_id": UUID(document_id),
                "workspace_id": UUID(workspace_id),
                "original_filename": original_filename,
                "content_text": content,
                "metadata": json.dumps(metadata) if metadata else None,
                "source_type": source_type,
                "primary_embedding_model": EMBEDDING_MODEL,
                "file_hash": file_hash,
            },
        )

    # ----------------------------
    # Insert new with dedupe
    # ----------------------------
    else:
        # If file_hash is missing, we cannot dedupe. We still insert.
        # (Best practice is to compute from raw bytes upstream and pass it in.)
        if file_hash is not None:
            # IMPORTANT: Because the recommended unique index is PARTIAL (WHERE file_hash IS NOT NULL),
            # the ON CONFLICT target must include the same predicate so Postgres can infer it.
            result = db.execute(
                text("""
                    INSERT INTO documents (
                        workspace_id,
                        original_filename,
                        content_text,
                        metadata,
                        source_type,
                        primary_embedding_model,
                        file_hash
                    )
                    VALUES (
                        :workspace_id,
                        :original_filename,
                        :content_text,
                        :metadata,
                        :source_type,
                        :primary_embedding_model,
                        :file_hash
                    )
                    ON CONFLICT (workspace_id, file_hash) WHERE file_hash IS NOT NULL
                    DO NOTHING
                    RETURNING id;
                """),
                {
                    "workspace_id": UUID(workspace_id),
                    "original_filename": original_filename,
                    "content_text": content,
                    "metadata": json.dumps(metadata) if metadata else None,
                    "source_type": source_type,
                    "primary_embedding_model": EMBEDDING_MODEL,
                    "file_hash": file_hash,
                },
            )
            doc_row = result.fetchone()

            if doc_row is None:
                # Duplicate: fetch existing id, skip embeddings.
                existing = db.execute(
                    text("""
                        SELECT id
                        FROM documents
                        WHERE workspace_id = :workspace_id
                          AND file_hash = :file_hash
                        LIMIT 1;
                    """),
                    {"workspace_id": UUID(workspace_id), "file_hash": file_hash},
                ).fetchone()

                existing_id = str(existing.id) if existing else None

                return {
                    "document_id": existing_id,
                    "num_chunks": 0,
                    "source_type": source_type,
                    "replaced": False,
                    "file_hash": file_hash,
                    "deduped": True,
                }

            document_id = str(doc_row.id)

        else:
            # No hash, do a plain insert.
            result = db.execute(
                text("""
                    INSERT INTO documents (
                        workspace_id,
                        original_filename,
                        content_text,
                        metadata,
                        source_type,
                        primary_embedding_model,
                        file_hash
                    )
                    VALUES (
                        :workspace_id,
                        :original_filename,
                        :content_text,
                        :metadata,
                        :source_type,
                        :primary_embedding_model,
                        :file_hash
                    )
                    RETURNING id;
                """),
                {
                    "workspace_id": UUID(workspace_id),
                    "original_filename": original_filename,
                    "content_text": content,
                    "metadata": json.dumps(metadata) if metadata else None,
                    "source_type": source_type,
                    "primary_embedding_model": EMBEDDING_MODEL,
                    "file_hash": file_hash,
                },
            )
            doc_row = result.fetchone()
            document_id = str(doc_row.id)

    # ----------------------------
    # DEBUG: read back what DB currently has for file_hash (within same transaction)
    # ----------------------------
    try:
        row = db.execute(
            text("""
                SELECT file_hash
                FROM documents
                WHERE id = :document_id
                  AND workspace_id = :workspace_id
                LIMIT 1;
            """),
            {"document_id": UUID(document_id), "workspace_id": UUID(workspace_id)},
        ).fetchone()
        logger.warning(
            "DEBUG DB file_hash readback (pre-commit): doc_id=%s file_hash_db=%r file_hash_arg=%r",
            document_id,
            (row[0] if row else None),
            file_hash,
        )
    except Exception as e:
        logger.exception("DEBUG failed to read back file_hash pre-commit: %s", repr(e))

    # ----------------------------
    # Embeddings insert (tenant_id is still required by embeddings table schema)
    # ----------------------------
    chunks = chunk_text(content)

    for idx, chunk in enumerate(chunks):
        embedding = embed_text(chunk)
        if isinstance(embedding, str):
            raise RuntimeError(f"embed_text returned a string instead of a vector for chunk {idx}")

        db.execute(
            text("""
                INSERT INTO embeddings (
                    tenant_id,
                    workspace_id,
                    document_id,
                    chunk_index,
                    chunk_text,
                    embedding,
                    embedding_model
                )
                VALUES (
                    :tenant_id,
                    :workspace_id,
                    :document_id,
                    :chunk_index,
                    :chunk_text,
                    :embedding,
                    :embedding_model
                );
            """),
            {
                "tenant_id": UUID(tenant_id),
                "workspace_id": UUID(workspace_id),
                "document_id": UUID(document_id),
                "chunk_index": idx,
                "chunk_text": chunk,
                "embedding": embedding,
                "embedding_model": EMBEDDING_MODEL,
            },
        )

    db.commit()

    # Optional: confirm after commit too (more definitive if you suspect later wipes)
    try:
        row2 = db.execute(
            text("""
                SELECT file_hash
                FROM documents
                WHERE id = :document_id
                  AND workspace_id = :workspace_id
                LIMIT 1;
            """),
            {"document_id": UUID(document_id), "workspace_id": UUID(workspace_id)},
        ).fetchone()
        logger.warning(
            "DEBUG DB file_hash readback (post-commit): doc_id=%s file_hash_db=%r file_hash_arg=%r",
            document_id,
            (row2[0] if row2 else None),
            file_hash,
        )
    except Exception as e:
        logger.exception("DEBUG failed to read back file_hash post-commit: %s", repr(e))

    return {
        "document_id": str(document_id),
        "num_chunks": len(chunks),
        "source_type": source_type,
        "replaced": replaced,
        "file_hash": file_hash,
        "deduped": False,
    }


def extract_text_from_pptx_bytes(content_bytes: bytes) -> str:
    if not content_bytes:
        return ""

    prs = Presentation(BytesIO(content_bytes))
    texts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_text_parts.append(t)

        if slide_text_parts:
            texts.append(f"[Slide {slide_idx}]")
            texts.append("\n".join(slide_text_parts))

    return "\n\n".join(texts).strip()


def extract_text_from_xlsx_bytes(content_bytes: bytes) -> str:
    if not content_bytes:
        return ""

    wb = load_workbook(BytesIO(content_bytes), data_only=True)
    lines: list[str] = []

    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) for cell in row if cell is not None]
            if row_values:
                lines.append("\t".join(row_values))
        lines.append("")

    return "\n".join(lines).strip()


def extract_text_from_docx_bytes(content_bytes: bytes) -> str:
    """
    Extract text from a .docx file provided as bytes.
    Note: .docx is a zip container; do NOT decode as UTF-8.
    """
    if not content_bytes:
        return ""

    doc = DocxDocument(io.BytesIO(content_bytes))
    paragraphs = [p.text for p in doc.paragraphs if (p.text or "").strip()]

    # Optional: include simple table text if present
    try:
        for table in doc.tables:
            for row in table.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                cells = [c for c in cells if c]
                if cells:
                    paragraphs.append(" | ".join(cells))
    except Exception:
        pass

    text_out = "\n".join(paragraphs)

    # Defensive: Postgres cannot store NUL bytes in text
    return text_out.replace("\x00", "")


def extract_text_from_upload(file: UploadFile) -> str:
    filename = file.filename or "unnamed"
    lower_name = filename.lower()

    content_bytes = file.file.read()

    if lower_name.endswith(".txt") or lower_name.endswith(".md"):
        return content_bytes.decode("utf-8", errors="ignore")

    if lower_name.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(content_bytes))
        pages_text = []
        for page in reader.pages:
            pages_text.append(page.extract_text() or "")
        return "\n".join(pages_text)

    if lower_name.endswith(".docx"):
        return extract_text_from_docx_bytes(content_bytes)

    if lower_name.endswith(".pptx"):
        return extract_text_from_pptx_bytes(content_bytes)

    if lower_name.endswith(".ppt"):
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported legacy PowerPoint file '{filename}'. "
                "Please resave the file as .pptx and upload again."
            ),
        )

    if lower_name.endswith(".xlsx"):
        return extract_text_from_xlsx_bytes(content_bytes)

    if lower_name.endswith(".xls"):
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported legacy Excel file '{filename}'. "
                "Please resave the file as .xlsx and upload again."
            ),
        )

    raise HTTPException(
        status_code=400,
        detail=(
            f"Unsupported file type for '{filename}'. "
            "Supported: .txt, .md, .pdf, .docx, .pptx, .xlsx"
        ),
    )
